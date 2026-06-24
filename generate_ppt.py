import sys
import os

# Include our local python packages path
local_packages = os.path.abspath('./python_packages')
if local_packages not in sys.path:
    sys.path.insert(0, local_packages)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9 (13.33" x 7.5")
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Use blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 1. Background (Full slide dark rectangle)
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 23, 42)  # slate-900 (#0f172a)
    bg.line.fill.background()  # No border
    
    # 2. Main Title Text Box
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.33), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = "Margin Slip Detector"
    p.font.name = "Arial"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p2 = tf.add_paragraph()
    p2.text = "Live P&L forecasting, governance alerts, and AI action playbooks built for Rocketlane"
    p2.font.name = "Arial"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(148, 163, 184)  # slate-400
    p2.space_before = Pt(8)

    # 3. Cards Data (Title, Description, Border Color, Accent Color, Left Position)
    cards_data = [
        {
            "title": "PORTFOLIO HUB",
            "desc": "Live margin forecasting across all active projects. Computes Estimate at Completion (EAC), billable ratios, and project profit margins in real time to capture slippage early.",
            "color": RGBColor(59, 130, 246),  # Blue
            "left": 1.0
        },
        {
            "title": "GOVERNANCE & PLAYBOOKS",
            "desc": "Automated risk tiering (Critical, High, Medium) paired with 4-point response playbooks for immediate margin protection. Standardizes organizational governance.",
            "color": RGBColor(245, 158, 11),  # Amber
            "left": 4.91
        },
        {
            "title": "NITRO AI BRIEFS",
            "desc": "Powered by Anthropic Claude to generate instant project health briefs, PM checklists, change-order recommendations, and draft stakeholder escalation emails with one click.",
            "color": RGBColor(139, 92, 246),  # Violet
            "left": 8.83
        }
    ]
    
    # 4. Generate Cards
    for c in cards_data:
        # Card Background Shape
        card = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(c["left"]), Inches(2.6), Inches(3.5), Inches(3.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(30, 41, 59)  # slate-800 (#1e293b)
        card.line.color.rgb = c["color"]
        card.line.width = Pt(2)
        
        # Card Text Box
        card_text_box = slide.shapes.add_textbox(
            Inches(c["left"] + 0.25), Inches(2.85), Inches(3.0), Inches(3.3)
        )
        ctf = card_text_box.text_frame
        ctf.word_wrap = True
        ctf.margin_left = ctf.margin_top = ctf.margin_right = ctf.margin_bottom = 0
        
        # Card Header
        cp1 = ctf.paragraphs[0]
        cp1.text = c["title"]
        cp1.font.name = "Arial"
        cp1.font.size = Pt(13)
        cp1.font.bold = True
        cp1.font.color.rgb = c["color"]
        cp1.space_after = Pt(14)
        
        # Card Description
        cp2 = ctf.add_paragraph()
        cp2.text = c["desc"]
        cp2.font.name = "Arial"
        cp2.font.size = Pt(12)
        cp2.font.color.rgb = RGBColor(203, 213, 225)  # slate-300
        cp2.line_spacing = 1.3
        
    # 5. Footer Text Box
    footer_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.7), Inches(11.33), Inches(0.4))
    ftf = footer_box.text_frame
    ftf.word_wrap = True
    ftf.margin_left = ftf.margin_top = ftf.margin_right = ftf.margin_bottom = 0
    
    fp = ftf.paragraphs[0]
    fp.text = "Built on Rocketlane RLI Framework × Anthropic Claude Sonnet 4.6"
    fp.font.name = "Arial"
    fp.font.size = Pt(10)
    fp.font.italic = True
    fp.font.color.rgb = RGBColor(100, 116, 139)  # slate-500
    
    # Save the presentation
    output_path = './presentation.pptx'
    prs.save(output_path)
    print(f"Presentation saved successfully at: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
